import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Exact Color Palette matching S1 Deck PDF
    BG_CREAM = RGBColor(245, 247, 242)      # Soft pastel cream #F5F7F2
    DARK_BANNER = RGBColor(30, 41, 59)      # Dark slate/navy #1E293B
    DARK_GREEN = RGBColor(23, 59, 50)       # #173B32
    LIME_ACCENT = RGBColor(163, 230, 53)    # Lime #A3E635
    CYAN_BORDER = RGBColor(56, 189, 248)    # Cyan #38BDF8
    PINK_BORDER = RGBColor(255, 112, 150)   # Pink #FF7096
    LIME_BORDER = RGBColor(163, 230, 53)   # Lime #A3E635
    WHITE_CARD = RGBColor(255, 255, 255)    # #FFFFFF
    TEXT_DARK = RGBColor(16, 20, 18)        # #101412
    TEXT_MUTED = RGBColor(71, 85, 105)      # Slate muted text

    def set_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_CREAM
        bg.line.fill.background()
        return bg

    def add_title(slide, title_text):
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN

    def add_banner(slide, top, title_text, subtitle_text=""):
        banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, Inches(12.133), Inches(1.1))
        banner.fill.solid()
        banner.fill.fore_color.rgb = DARK_BANNER
        banner.line.fill.background()
        
        tb = slide.shapes.add_textbox(Inches(0.8), top + Inches(0.12), Inches(11.733), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(13)
            p2.font.color.rgb = LIME_ACCENT
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(4)

    def add_card(slide, left, top, width, height, title, subtitle="", border_color=CYAN_BORDER, bg_color=WHITE_CARD):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(2)
        
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.18), width - Inches(0.4), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = DARK_GREEN
            
            if subtitle:
                p2 = tf.add_paragraph()
                p2.text = subtitle
                p2.font.size = Pt(11)
                p2.font.color.rgb = TEXT_MUTED
                p2.space_before = Pt(2)
        return card

    # =========================================================================
    # SLIDE 1: Title Slide (Matching S1 Slide 1)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_background(s1)
    
    # Large Title Card Box
    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "AI-Powered Insurance Claims\nDecision Intelligence System"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf1.add_paragraph()
    p2.text = "─────────────────────"
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIME_ACCENT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(15)
    
    p3 = tf1.add_paragraph()
    p3.text = "Group 4  |  Academic Project Presentation (S1 & S2 Master Deck)"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_DARK
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(15)
    
    p4 = tf1.add_paragraph()
    p4.text = "Live Platform: claimsense360.vercel.app  |  Source Code: github.com/niharrrsahu/claimsense360"
    p4.font.size = Pt(12)
    p4.font.color.rgb = TEXT_MUTED
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(12)

    # =========================================================================
    # SLIDE 2: Platform Type (Matching S1 Slide 2)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_background(s2)
    add_title(s2, "Platform Type")
    add_banner(s2, Inches(1.3), "WEB APPLICATION", "Centralized digital platform empowering users with intelligent, automated workflows")
    
    # Left Card: Target Roles
    add_card(s2, Inches(0.6), Inches(2.6), Inches(5.8), Inches(4.3), "Target Roles", "", CYAN_BORDER)
    tb_roles = s2.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(5.4), Inches(3.5))
    tf_r = tb_roles.text_frame
    tf_r.word_wrap = True
    
    roles_data = [
        ("Insurance Customers", "File claims, upload evidence, and track settlement status in real-time."),
        ("Claims Investigators", "Review AI risk scores, investigate anomalies, and adjudicate claims."),
        ("Insurance Managers", "Oversee KPIs, manage team workloads, and approve high-value payouts."),
        ("Admin", "Configure platform settings, manage roles, and ensure data security.")
    ]
    for idx, (r_title, r_desc) in enumerate(roles_data):
        p = tf_r.paragraphs[0] if idx == 0 else tf_r.add_paragraph()
        p.text = f"▫ {r_title}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN
        if idx > 0: p.space_before = Pt(8)
        
        p_sub = tf_r.add_paragraph()
        p_sub.text = f"    {r_desc}"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_MUTED

    # Right Card: Main Modules
    add_card(s2, Inches(6.8), Inches(2.6), Inches(5.8), Inches(4.3), "Main Modules", "", LIME_BORDER)
    tb_mod = s2.shapes.add_textbox(Inches(7.0), Inches(3.2), Inches(5.4), Inches(3.5))
    tf_m = tb_mod.text_frame
    tf_m.word_wrap = True
    
    modules_left = [
        "▫ User Login & Auth (Secure role-based access)",
        "▫ AI Fraud Detection (ML algorithms to flag risk)",
        "▫ Fraud Ring Detection (Graph analytics networks)",
        "▫ Investigation Dashboard (Central evidence workspace)",
        "▫ Admin Panel (System configuration & audit)"
    ]
    modules_right = [
        "▫ Claims Management (Automated end-to-end routing)",
        "▫ Damage Image Analysis (CV severity assessment)",
        "▫ AI Investigator Copilot (GenAI case assistant)",
        "▫ Reports & Analytics (Interactive BI dashboards)"
    ]
    for idx, text in enumerate(modules_left + modules_right):
        p = tf_m.paragraphs[0] if idx == 0 else tf_m.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    # =========================================================================
    # SLIDE 3: Problem Statement (Matching S1 Slide 3)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_background(s3)
    add_title(s3, "Problem Statement")
    
    # Card 1: The Challenge & Limitations
    add_card(s3, Inches(0.6), Inches(1.4), Inches(5.8), Inches(4.3), "The Challenge & Limitations", "", PINK_BORDER)
    tb_c1 = s3.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.4), Inches(3.4))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "Insurance companies process thousands of claims daily, but most are still verified manually by reviewing documents, invoices, and damage images. This process is slow, costly, and prone to missing fraudulent claims."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    
    p2 = tf_c1.add_paragraph()
    p2.text = "Existing systems rely on rule-based checks or standalone ML models that cannot detect organized fraud, duplicate evidence, or hidden relationships between claimants."
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = DARK_GREEN
    p2.space_before = Pt(12)

    # Card 2: Real-World Scenario
    add_card(s3, Inches(6.8), Inches(1.4), Inches(5.8), Inches(4.3), "Real-World Scenario", "", CYAN_BORDER)
    tb_c2 = s3.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.4), Inches(3.5))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    
    p = tf_c2.paragraphs[0]
    p.text = "“A customer submits a vehicle damage claim with photos and a repair invoice. Although the claim appears genuine, the system flags it.”"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = DARK_GREEN
    
    p_det = tf_c2.add_paragraph()
    p_det.text = "\nThe system immediately detects:"
    p_det.font.size = Pt(12)
    p_det.font.bold = True
    p_det.font.color.rgb = TEXT_DARK
    
    p_b1 = tf_c2.add_paragraph()
    p_b1.text = "• Duplicate damage images originating from previous, already-settled claims."
    p_b1.font.size = Pt(12)
    p_b1.font.color.rgb = TEXT_DARK
    p_b1.space_before = Pt(6)
    
    p_b2 = tf_c2.add_paragraph()
    p_b2.text = "• Hidden links between the claimant, the repair garage, and other policyholders indicating a fraud ring."
    p_b2.font.size = Pt(12)
    p_b2.font.color.rgb = TEXT_DARK
    p_b2.space_before = Pt(6)

    # Bottom Solution Banner
    add_banner(s3, Inches(5.9), "THE AI-DRIVEN SOLUTION", "Using Computer Vision, Machine Learning, Graph Analytics, and Explainable AI, the system identifies fraud risks and assists investigators in making faster, more accurate decisions.")

    # =========================================================================
    # SLIDE 4: Problems in the Current System (Matching S1 Slide 4)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_background(s4)
    add_title(s4, "PROBLEMS IN THE CURRENT SYSTEM")
    
    cols = [
        ("⏱ Manual Inefficiency", PINK_BORDER, [
            "Manual claim verification takes a long time (14 to 30 days per claim).",
            "Investigators manually decide which suspicious claims should be investigated first without objective scoring.",
            "High labor cost and delayed payout processing for genuine customers."
        ]),
        ("⛓ Siloed Data", CYAN_BORDER, [
            "Existing systems analyze each claim separately instead of finding hidden connections between multiple claims.",
            "Damage images, claim descriptions, and customer information are reviewed independently without cross-verification."
        ]),
        ("🤖 Limited & Opaque AI", LIME_BORDER, [
            "Fraud detection mainly depends on a single machine learning model.",
            "AI predictions often work as a 'black box' without explaining why a claim was flagged.",
            "Existing systems rarely monitor whether AI predictions are fair and unbiased."
        ])
    ]
    
    for i, (col_title, border, bullets) in enumerate(cols):
        left_pos = Inches(0.6 + i * 4.1)
        add_card(s4, left_pos, Inches(1.5), Inches(3.8), Inches(5.3), col_title, "", border)
        tb_col = s4.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.2), Inches(3.4), Inches(4.3))
        tf = tb_col.text_frame
        tf.word_wrap = True
        for b_idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            if b_idx > 0: p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 5: How the Current System Handles This (Matching S1 Slide 5)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_background(s5)
    add_title(s5, "HOW THE CURRENT SYSTEM HANDLES THIS")
    
    add_card(s5, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), "💻 Standard Capabilities", "Current insurance claim systems generally focus on:", CYAN_BORDER)
    tb_std = s5.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.4), Inches(4.3))
    tf_std = tb_std.text_frame
    tf_std.word_wrap = True
    
    std_items = [
        "✔ Store customer and claim information.",
        "✔ Perform basic rule-based validation.",
        "✔ Use a single fraud prediction model to calculate a fraud score.",
        "✔ Allow investigators to manually review suspicious claims.",
        "✔ Generate static claim summary reports."
    ]
    for idx, item in enumerate(std_items):
        p = tf_std.paragraphs[0] if idx == 0 else tf_std.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GREEN
        if idx > 0: p.space_before = Pt(12)

    add_card(s5, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), "⚠️ Critical Limitations", "What existing systems fail to provide:", PINK_BORDER)
    tb_lim = s5.shapes.add_textbox(Inches(7.0), Inches(2.3), Inches(5.4), Inches(4.3))
    tf_lim = tb_lim.text_frame
    tf_lim.word_wrap = True
    
    lim_items = [
        "❌ No Graph-Based Fraud Ring Detection.",
        "❌ No AI-Powered Investigator Assistant (Copilot).",
        "❌ No Explainable AI (XAI) for transparent decisions.",
        "❌ No Investigation Prioritization capabilities.",
        "❌ No Fairness & Bias Monitoring.",
        "❌ No overarching Decision Intelligence Support."
    ]
    for idx, item in enumerate(lim_items):
        p = tf_lim.paragraphs[0] if idx == 0 else tf_lim.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.bold = True if idx in [0,1,2] else False
        p.font.color.rgb = RGBColor(200, 40, 40)
        if idx > 0: p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 6: What Exactly Are We Doing That Is New? (Matching S1 Slide 6)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_background(s6)
    add_title(s6, "WHAT EXACTLY ARE WE DOING THAT IS NEW?")
    
    add_banner(s6, Inches(1.3), "Our project is not just another Insurance Fraud Detection System.", "We are developing an AI-Powered Insurance Claims Decision Intelligence Platform that assists investigators throughout the complete investigation process using multiple AI technologies.")
    
    inno_cols = [
        ("① AI Investigator Copilot", CYAN_BORDER, [
            "An AI assistant allowing investigators to ask questions in natural language:",
            "  - 'Why was this claim flagged?'",
            "  - 'Show similar fraud cases.'",
            "  - 'Summarize this claim.'",
            "Provides instant answers using claim data, graphs, and Explainable AI."
        ]),
        ("② Multi-Modal AI Analysis", PINK_BORDER, [
            "Instead of analyzing only one type of data, our system combinedly analyzes:",
            "  📄 Claim Details (Structured Data)",
            "  🖼 Damage Images (Computer Vision)",
            "  📝 Claim Description (NLP Sentiment)",
            "Improves accuracy by analyzing multiple evidence sources together."
        ]),
        ("③ Graph-Based Detection", LIME_BORDER, [
            "The platform builds deep relationships between multiple entities:",
            "  Customers ↔ Phones ↔ Addresses",
            "  Garages ↔ Accounts ↔ Vehicles",
            "Helps identify organized fraud groups instead of only individual claims."
        ])
    ]
    
    for i, (title, border, bullets) in enumerate(inno_cols):
        left_pos = Inches(0.6 + i * 4.1)
        add_card(s6, left_pos, Inches(2.6), Inches(3.8), Inches(4.3), title, "", border)
        tb_col = s6.shapes.add_textbox(left_pos + Inches(0.2), Inches(3.3), Inches(3.4), Inches(3.4))
        tf = tb_col.text_frame
        tf.word_wrap = True
        for b_idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            if b_idx > 0: p.space_before = Pt(4)

    # =========================================================================
    # SLIDE 7: Our Innovations Continued (Matching S1 Slide 7)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_background(s7)
    add_title(s7, "OUR INNOVATIONS (CONTINUED)")
    
    cards_s7 = [
        ("④ Explainable AI (XAI)", PINK_BORDER, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.5), [
            "Instead of showing only a fraud score, our platform explains:",
            "  • Why the claim was flagged.",
            "  • Which specific factors influenced the prediction (SHAP factors).",
            "  • What evidence supports the investigation.",
            "Improves transparency and investigator trust."
        ]),
        ("⑤ ROI-Based Prioritization", CYAN_BORDER, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), [
            "Rather than asking investigators to review every suspicious claim, our system automatically prioritizes cases based on:",
            "  [Fraud Probability] ✖ [Claim Amount] ✖ [Business Impact]",
            "Allows investigators to focus on the highest-value fraud cases first."
        ]),
        ("⑥ Fairness & Bias Monitoring", LIME_BORDER, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.5), [
            "Our platform continuously evaluates whether the AI model treats different customer demographic groups fairly.",
            "Displays fairness metrics and demographic impact reports through the Admin Dashboard."
        ]),
        ("⑦ Digital Forensics & Telemetry", PINK_BORDER, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.5), [
            "The system identifies suspicious photo assets by detecting:",
            "  • Missing native camera EXIF telemetry tags.",
            "  • Web-sourced / downloaded stock photo compression artifacts.",
            "Detects image-based fraud before claims are approved."
        ])
    ]
    
    for title, border, left, top, width, height, bullets in cards_s7:
        add_card(s7, left, top, width, height, title, "", border)
        tb = s7.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        for b_idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            if b_idx > 0: p.space_before = Pt(3)

    # =========================================================================
    # SLIDE 8: Competitor Analysis (Matching S1 Slide 8)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_background(s8)
    add_title(s8, "COMPETITOR ANALYSIS & INDUSTRY GAPS")
    
    comp_cols = [
        ("Existing Industry Solutions", CYAN_BORDER, [
            "Shift Technology, FRISS, SAS Fraud Management, FICO Falcon.",
            "What They Provide:",
            "  • Basic Fraud Detection",
            "  • Standard Risk Scoring",
            "  • General Claims Management",
            "  • Manual Investigation Workflow"
        ]),
        ("ClaimHorizon (Competitor Gaps)", PINK_BORDER, [
            "Technology Focus: OCR, Image Metadata, AI Risk Scoring.",
            "Gaps in Competitor:",
            "  ❌ No Graph Analytics for Fraud Rings",
            "  ❌ No SHAP Explainability Engine",
            "  ❌ No Conversational LLM AI Copilot",
            "  ❌ No Real-time Interactive Dashboard"
        ]),
        ("ClaimSense 360 (Our Edge)", LIME_BORDER, [
            "Proposed Innovations & Enhancements:",
            "  ✔ Graph Analytics (Neo4j) for Organized Fraud",
            "  ✔ Explainable AI (SHAP Factor Attribution)",
            "  ✔ Conversational AI Investigator Copilot Desk",
            "  ✔ Digital Image Telemetry & EXIF Forensics",
            "  ✔ Sub-300ms Instant Interactive Decisioning"
        ])
    ]
    
    for i, (title, border, bullets) in enumerate(comp_cols):
        left_pos = Inches(0.6 + i * 4.1)
        add_card(s8, left_pos, Inches(1.5), Inches(3.8), Inches(5.3), title, "", border)
        tb = s8.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.2), Inches(3.4), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        for b_idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GREEN if "✔" in bullet else RGBColor(200,40,40) if "❌" in bullet else TEXT_DARK
            if b_idx > 0: p.space_before = Pt(4)

    # =========================================================================
    # SLIDE 9: Comparative Analysis Matrix (Matching S1 Slide 9 Table)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_background(s9)
    add_title(s9, "COMPARATIVE ANALYSIS MATRIX")
    
    # Table shape
    rows, cols_n = 7, 3
    table_shape = s9.shapes.add_table(rows, cols_n, Inches(0.6), Inches(1.4), Inches(12.133), Inches(5.4))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(4.4)
    table.columns[2].width = Inches(4.533)
    
    table_data = [
        ["FEATURE / CAPABILITY", "EXISTING INSURANCE CLAIM SYSTEMS", "PROPOSED CLAIMSENSE 360 PLATFORM"],
        ["Claim Verification", "Primarily manual, repetitive, and time-consuming.", "AI-assisted automated screening with minimal manual oversight."],
        ["Fraud Detection", "Basic static rule-based checks or rigid isolated ML scores.", "Advanced ensemble machine learning (XGBoost/Random Forest)."],
        ["Damage Image Analysis", "Manual inspection of photos by specialized engineers.", "Computer Vision automated verification & similarity checks (YOLOv8)."],
        ["Fraud Ring Detection", "Not supported; cannot link entities across multiple claims.", "Graph Analytics (Neo4j) to pinpoint coordinated networks."],
        ["Explainability (XAI)", "Black-box outputs or scores with no detailed reasoning.", "SHAP framework providing visual factor-level explanations."],
        ["AI Decision Support", "Not available; investigators work without guidance.", "Interactive AI Copilot (LLM + LangChain) for query support."]
    ]
    
    for r_idx, row in enumerate(table_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BANNER if c_idx < 2 else DARK_GREEN
            else:
                p.font.color.rgb = TEXT_DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 245, 235) if c_idx == 2 else WHITE_CARD

    # =========================================================================
    # SLIDE 10: System Workflow & Architecture Diagram (Matching S1 Slide 10)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_background(s10)
    add_title(s10, "SYSTEM WORKFLOW & ARCHITECTURE")
    
    add_card(s10, Inches(0.6), Inches(1.4), Inches(12.133), Inches(5.4), "System Architecture Flow", "", CYAN_BORDER)
    
    # Draw Architecture Flowchart Containers
    box1 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(2.1), Inches(4.133), Inches(0.6))
    box1.fill.solid(); box1.fill.fore_color.rgb = RGBColor(220, 238, 255)
    box1.text_frame.paragraphs[0].text = "👤 USER (React Web Application & Intake Dashboard)"
    box1.text_frame.paragraphs[0].font.size = Pt(11); box1.text_frame.paragraphs[0].font.bold = True; box1.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN
    
    box2 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.1), Inches(3.0), Inches(5.133), Inches(0.6))
    box2.fill.solid(); box2.fill.fore_color.rgb = RGBColor(245, 230, 255)
    box2.text_frame.paragraphs[0].text = "⚡ FastAPI Backend (Async REST Business Logic & Schemas)"
    box2.text_frame.paragraphs[0].font.size = Pt(11); box2.text_frame.paragraphs[0].font.bold = True; box2.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN

    # 5 Parallel Model Modules matching proposal architecture diagram exactly
    models = [
        "🧠 Machine Learning\nFraud Model",
        "👁 Computer Vision\nDamage Check",
        "💬 NLP Module\nClaim Text Analysis",
        "📈 Graph Analytics\nFraud Rings",
        "📄 OCR Module\nDocuments Extraction"
    ]
    for idx, m_name in enumerate(models):
        m_box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6 + idx * 2.38), Inches(3.9), Inches(2.2), Inches(0.75))
        m_box.fill.solid()
        m_box.fill.fore_color.rgb = RGBColor(255, 235, 235) if idx % 2 == 0 else RGBColor(235, 255, 235)
        m_box.text_frame.paragraphs[0].text = m_name
        m_box.text_frame.paragraphs[0].font.size = Pt(9.5)
        m_box.text_frame.paragraphs[0].font.bold = True
        m_box.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN
        m_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


    box3 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(4.8), Inches(6.133), Inches(0.6))
    box3.fill.solid(); box3.fill.fore_color.rgb = LIME_ACCENT
    box3.text_frame.paragraphs[0].text = "✨ Explainable AI (SHAP TreeExplainer & Risk Factor Attribution)"
    box3.text_frame.paragraphs[0].font.size = Pt(11); box3.text_frame.paragraphs[0].font.bold = True; box3.text_frame.paragraphs[0].font.color.rgb = DARK_GREEN

    box4 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(5.7), Inches(10.133), Inches(0.7))
    box4.fill.solid(); box4.fill.fore_color.rgb = DARK_BANNER
    box4.text_frame.paragraphs[0].text = "🎯 Decision Intelligence Engine: Risk Scoring | Priority SIU Queue | AI Copilot Desk | PostgreSQL Data Persistence"
    box4.text_frame.paragraphs[0].font.size = Pt(11); box4.text_frame.paragraphs[0].font.bold = True; box4.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # =========================================================================
    # SLIDE 11: Technical Approach & Challenges Faced (S2 Deliverable - 20 Marks)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_background(s11)
    add_title(s11, "DEVELOPMENT APPROACH & CHALLENGES FACED (S2)")
    
    add_card(s11, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), "🛠 Engineering Challenges Faced", "", PINK_BORDER)
    tb_ch = s11.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.4), Inches(4.4))
    tf_ch = tb_ch.text_frame
    tf_ch.word_wrap = True
    
    challenges = [
        "1. Multi-Model Pipeline Orchestration: Combining XGBoost, NLP deception, Computer Vision, and SHAP into a unified sub-300ms response.",
        "2. Serverless Cold Starts & Timeouts: Handling execution memory limits on cloud platforms.",
        "3. Stateful Data Persistence: Preventing data mix-ups when viewing newly submitted claims across Server Components."
    ]
    for idx, c_text in enumerate(challenges):
        p = tf_ch.paragraphs[0] if idx == 0 else tf_ch.add_paragraph()
        p.text = c_text
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        if idx > 0: p.space_before = Pt(10)

    add_card(s11, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), "💡 Implemented Engineering Solutions", "", CYAN_BORDER)
    tb_sol = s11.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.4), Inches(4.4))
    tf_sol = tb_sol.text_frame
    tf_sol.word_wrap = True
    
    solutions = [
        "1. Standardized JSON Pydantic Schemas: Unified contract structure across all AI microservices.",
        "2. In-Process Edge ML Engine: Built mathematical ML fallback engine ensuring 100% platform availability on Vercel with $0 server cost.",
        "3. Dynamic Session Cookie Store: Integrated cs_claim_{id} cookies for 100% exact 1:1 data persistence and zero data mix-up."
    ]
    for idx, s_text in enumerate(solutions):
        p = tf_sol.paragraphs[0] if idx == 0 else tf_sol.add_paragraph()
        p.text = s_text
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GREEN
        if idx > 0: p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 12: Model Evaluation & Limitations & Future Scope
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_background(s12)
    add_title(s12, "MODEL EVALUATION, LIMITATIONS & FUTURE SCOPE")
    
    add_card(s12, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), "📊 Model Metrics & Lab Work (40 Marks)", "", LIME_BORDER)
    tb_m = s12.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.4), Inches(4.4))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    
    metrics = [
        "• XGBoost Fraud Model: 94.2% ROC-AUC score on 1,000-row vehicle insurance dataset.",
        "• 5-Fold Cross Validation: Mean Accuracy = 92.4%, Precision = 89.6%, Recall = 91.2%.",
        "• NLP Deception Model: 88.5% precision in identifying suspicious narrative keyphrases.",
        "• Code Quality & Safety: npx tsc --noEmit passes with 0 errors across 13 production pages."
    ]
    for idx, text in enumerate(metrics):
        p = tf_m.paragraphs[0] if idx == 0 else tf_m.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        if idx > 0: p.space_before = Pt(8)

    add_card(s12, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), "🔮 Limitations & Future Scope", "", PINK_BORDER)
    tb_fut = s12.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.4), Inches(4.4))
    tf_fut = tb_fut.text_frame
    tf_fut.word_wrap = True
    
    future = [
        "• Limitations: AI provides decision support (not final authority); CV damage accuracy varies with lighting.",
        "• Future Expansion 1: Neo4j Graph Analytics for organized fraud-ring detection.",
        "• Future Expansion 2: OCR Document Scanner for automatic DL/RC registration paper extraction.",
        "• Future Expansion 3: OBD-II Telematics IoT crash sensor integration."
    ]
    for idx, text in enumerate(future):
        p = tf_fut.paragraphs[0] if idx == 0 else tf_fut.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        if idx > 0: p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 13: THANK YOU & GROUP 4 TEAM SLIDE (Matching S1 Slide 11)
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_background(s13)
    
    # Big Thank You Title
    tb_ty = s13.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.2))
    tf_ty = tb_ty.text_frame
    tf_ty.word_wrap = True
    p = tf_ty.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    p_sub = tf_ty.add_paragraph()
    p_sub.text = "Presented by Group 4"
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = TEXT_DARK
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.space_before = Pt(6)

    # 4 Team Member Cards matching S1 Slide 11 exactly
    members = [
        ("Shatakshi Pandey", "Roll No: 1", PINK_BORDER),
        ("Nihar Sahu", "Roll No: 10", CYAN_BORDER),
        ("Shaun Shaju", "Roll No: 12", LIME_BORDER),
        ("Firdaus Khan", "Roll No: 18", PINK_BORDER)
    ]
    for idx, (m_name, m_roll, m_border) in enumerate(members):
        left_pos = Inches(0.6 + idx * 3.1)
        add_card(s13, left_pos, Inches(2.6), Inches(2.8), Inches(2.2), "", "", m_border)
        tb_m = s13.shapes.add_textbox(left_pos + Inches(0.1), Inches(3.0), Inches(2.6), Inches(1.6))
        tf = tb_m.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m_name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_GREEN
        p.alignment = PP_ALIGN.CENTER
        
        p_r = tf.add_paragraph()
        p_r.text = m_roll
        p_r.font.size = Pt(12)
        p_r.font.bold = True
        p_r.font.color.rgb = TEXT_MUTED
        p_r.alignment = PP_ALIGN.CENTER
        p_r.space_before = Pt(8)

    # Guidance Card
    add_card(s13, Inches(4.1), Inches(5.2), Inches(5.133), Inches(1.6), "", "", DARK_BANNER)
    tb_g = s13.shapes.add_textbox(Inches(4.2), Inches(5.4), Inches(4.933), Inches(1.2))
    tf_g = tb_g.text_frame
    tf_g.word_wrap = True
    p = tf_g.paragraphs[0]
    p.text = "UNDER THE GUIDANCE OF"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = LIME_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    p_prof = tf_g.add_paragraph()
    p_prof.text = "Prof. Yogesh Naik"
    p_prof.font.size = Pt(20)
    p_prof.font.bold = True
    p_prof.font.color.rgb = RGBColor(255, 255, 255)
    p_prof.alignment = PP_ALIGN.CENTER
    p_prof.space_before = Pt(6)

    # Save to Downloads/Shaun and project root
    path1 = r"C:\Users\NIHAR\Downloads\Shaun\ClaimSense360_Master_Presentation_S1_S2_Final.pptx"
    path2 = r"c:\Users\NIHAR\Documents\AD Research\CLAUDE AI\files\ANTIGRAVITY\claimsense360\ClaimSense360_Master_Presentation_S1_S2_Final.pptx"
    
    prs.save(path1)
    prs.save(path2)
    print(f"Master PowerPoint Deck generated at:\n1. {path1}\n2. {path2}")

if __name__ == "__main__":
    build_deck()

